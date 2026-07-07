import sys
import os
import openpyxl
from pptx import Presentation

# Set stdout encoding just in case, but write directly to file in UTF-8
def extract_pptx_text(pptx_path, out_file):
    out_file.write(f"\n=========================================\nEXTRACTING FROM: {os.path.basename(pptx_path)}\n=========================================\n\n")
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        out_file.write(f"--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                out_file.write(shape.text.strip() + "\n")
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    row_text = []
                    for col_idx, cell in enumerate(row.cells):
                        row_text.append(cell.text.strip())
                    out_file.write(f"Table Row {row_idx+1}: " + " | ".join(row_text) + "\n")
        out_file.write("\n")

def extract_xlsx_info(xlsx_path, out_file):
    out_file.write(f"\n=========================================\nEXTRACTING FROM: {os.path.basename(xlsx_path)}\n=========================================\n\n")
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    for sheet_name in wb.sheetnames:
        out_file.write(f"--- Sheet: {sheet_name} ---\n")
        sheet = wb[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                out_file.write(row_str + "\n")
        out_file.write("\n")

if __name__ == "__main__":
    pptx1 = "Universal Sompo Evidence Dashboard Plan.pptx"
    pptx2 = "Universal_Sompo_AI_Claim_Evidence_Finder_Tech_Plan.pptx"
    xlsx1 = "Universal_Sompo_4_Week_Build_and_GPU_Cost_Plan.xlsx"
    
    with open("extracted_context.txt", "w", encoding="utf-8") as out:
        if os.path.exists(pptx1):
            extract_pptx_text(pptx1, out)
        else:
            out.write(f"File not found: {pptx1}\n")
            
        if os.path.exists(pptx2):
            extract_pptx_text(pptx2, out)
        else:
            out.write(f"File not found: {pptx2}\n")
            
        if os.path.exists(xlsx1):
            extract_xlsx_info(xlsx1, out)
        else:
            out.write(f"File not found: {xlsx1}\n")
    print("Done extracting!")
